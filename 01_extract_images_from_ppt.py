from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import argparse
import cv2
import os

ap = argparse.ArgumentParser()
ap.add_argument("-p", "--ppt", type=str, default="Numpy and Pandas2.pptx",
    help="path to the ppt file")
ap.add_argument("-e", "--export_image_folder", type=str, default="images_extracted_from_ppt",
    help="path to extract images from ppt")
args = vars(ap.parse_args())

filename = args["ppt"]
export_folder_name = args["export_image_folder"]
file_names = []
text_from_image = {}

if not os.path.exists(export_folder_name):
    os.makedirs(export_folder_name)

# Function to check if images are the same
def check_image_if_same(imageA, imageB):
    imgA = cv2.imread(imageA)
    imgB = cv2.imread(imageB)
    if imgA.shape == imgB.shape:
        difference = cv2.subtract(imgA, imgB)
        b, g, r = cv2.split(difference)
        if cv2.countNonZero(b) == 0 and cv2.countNonZero(g) == 0 and cv2.countNonZero(r) == 0:
            return True
    return False

# Iterate the power point and extract the images
def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape


# Extract all the images and write in a folder
for index, picture in enumerate(iter_picture_shapes(Presentation(filename))):
    image = picture.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    # image_filename = 'image.%s' % image.ext
    image_filename = f"{export_folder_name}/{index}.{image.ext}"
    file_names.append(image_filename)
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)
    print(f"[INFO] Extracting image {index}.png ...")

# Remvoing Duplicate images
for index, i in enumerate(file_names):
    for j in file_names[index+1:]:
        if check_image_if_same(i, j):
            print(f"[INFO] Removing file {i}")
            os.remove(i)
            break
print("[INFO] Done!")